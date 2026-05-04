from __future__ import annotations

from collections.abc import Callable, Sequence
from pathlib import Path

import pytest


@pytest.fixture
def pdf_source_factory(tmp_path: Path) -> Callable[[Sequence[str], str], Path]:
    def factory(pages: Sequence[str], file_name: str = "sample.pdf") -> Path:
        path = tmp_path / file_name
        _write_text_layer_pdf(path, pages)
        return path

    return factory


@pytest.fixture
def pptx_source_factory(tmp_path: Path) -> Callable[[str], Path]:
    def factory(file_name: str = "sample.pptx") -> Path:
        pytest.importorskip("pptx")
        from pptx import Presentation
        from pptx.util import Inches

        path = tmp_path / file_name
        presentation = Presentation()

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Launch Checklist"
        slide.placeholders[1].text = "Approve release window and rollback contact."

        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Risk Review"
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(1))
        textbox.text_frame.text = "Escalate invoice anomalies before launch."
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.6), Inches(5), Inches(1))
        table = table_shape.table
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Target"
        table.cell(1, 0).text = "SLA-42"
        table.cell(1, 1).text = "99 percent"

        presentation.save(path)
        return path

    return factory


def _write_text_layer_pdf(path: Path, pages: Sequence[str]) -> None:
    objects: dict[int, bytes] = {
        1: b"<< /Type /Catalog /Pages 2 0 R >>",
        3: b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    }
    page_object_numbers: list[int] = []
    next_object_number = 4

    for page_text in pages:
        page_object_number = next_object_number
        content_object_number = next_object_number + 1
        next_object_number += 2
        page_object_numbers.append(page_object_number)

        stream = f"BT /F1 12 Tf 72 720 Td ({_pdf_string(page_text)}) Tj ET".encode("latin-1")
        objects[page_object_number] = (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 3 0 R >> >> "
            + f"/Contents {content_object_number} 0 R >>".encode("ascii")
        )
        objects[content_object_number] = (
            f"<< /Length {len(stream)} >>\nstream\n".encode("ascii")
            + stream
            + b"\nendstream"
        )

    kids = " ".join(f"{object_number} 0 R" for object_number in page_object_numbers)
    objects[2] = f"<< /Type /Pages /Kids [{kids}] /Count {len(pages)} >>".encode("ascii")

    max_object_number = max(objects)
    output_parts = [b"%PDF-1.4\n"]
    offsets = [0] * (max_object_number + 1)

    for object_number in range(1, max_object_number + 1):
        offsets[object_number] = sum(len(part) for part in output_parts)
        output_parts.append(f"{object_number} 0 obj\n".encode("ascii"))
        output_parts.append(objects[object_number])
        output_parts.append(b"\nendobj\n")

    start_xref = sum(len(part) for part in output_parts)
    output_parts.append(f"xref\n0 {max_object_number + 1}\n".encode("ascii"))
    output_parts.append(b"0000000000 65535 f \n")
    for object_number in range(1, max_object_number + 1):
        output_parts.append(f"{offsets[object_number]:010d} 00000 n \n".encode("ascii"))
    output_parts.append(
        f"trailer\n<< /Size {max_object_number + 1} /Root 1 0 R >>\n"
        f"startxref\n{start_xref}\n%%EOF\n".encode("ascii")
    )

    path.write_bytes(b"".join(output_parts))


def _pdf_string(value: str) -> str:
    return value.replace("\\", "\\\\").replace("(", r"\(").replace(")", r"\)")
